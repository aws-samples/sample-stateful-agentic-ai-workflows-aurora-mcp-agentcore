"""
Build a 2-slide PPTX deck explaining the ClickShop -> Meridian pivot.

Run: python scripts/build_pivot_deck.py
Output: docs/meridian-pivot.pptx (16:9, ready to open in Keynote / PowerPoint)
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt


# Meridian palette
INK = RGBColor(0x0E, 0x1B, 0x2C)
INK_SOFT = RGBColor(0x4A, 0x55, 0x68)
PAPER = RGBColor(0xFA, 0xF7, 0xF1)
ACCENT = RGBColor(0x1E, 0x3A, 0x8A)
ACCENT_SOFT = RGBColor(0xE8, 0xEE, 0xFA)
RULE = RGBColor(0xCB, 0xD2, 0xDC)
GREEN = RGBColor(0x16, 0x82, 0x4D)
MUTED = RGBColor(0x8B, 0x95, 0xA5)


PHASE_COLORS = [
    RGBColor(0x1E, 0x3A, 0x8A),  # 1 - blue
    RGBColor(0x0E, 0x74, 0x90),  # 2 - teal
    RGBColor(0x9A, 0x3F, 0x12),  # 3 - rust
    RGBColor(0x6B, 0x21, 0xA8),  # 4 - violet
    RGBColor(0x15, 0x6B, 0x3D),  # 5 - green
]


def _set_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _set_text(
    frame,
    text: str,
    *,
    size: int,
    bold: bool = False,
    color: RGBColor = INK,
    align=PP_ALIGN.LEFT,
    font: str = "Helvetica Neue",
) -> None:
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.alignment = align
    p.text = ""
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_text(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    size: int,
    bold: bool = False,
    color: RGBColor = INK,
    align=PP_ALIGN.LEFT,
    font: str = "Helvetica Neue",
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    _set_text(tf, text, size=size, bold=bold, color=color, align=align, font=font)
    return box


def _add_rect(slide, left: float, top: float, width: float, height: float, color: RGBColor):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    _set_fill(shape, color)
    return shape


def _add_rule(slide, left: float, top: float, width: float, color: RGBColor = RULE):
    return _add_rect(slide, left, top, width, 0.012, color)


def _slide_chrome(slide, *, slide_num: int, total: int, eyebrow: str) -> None:
    # Top-left eyebrow
    _add_text(
        slide, eyebrow, 0.5, 0.4, 8.0, 0.3,
        size=10, bold=True, color=ACCENT, font="Helvetica Neue",
    )
    # Top-right slide counter
    _add_text(
        slide, f"{slide_num} / {total}", 12.4, 0.4, 0.6, 0.3,
        size=10, color=MUTED, align=PP_ALIGN.RIGHT,
    )
    # Bottom rule + footer
    _add_rule(slide, 0.5, 7.05, 12.5)
    _add_text(
        slide, "Meridian — DAT309 workshop demo", 0.5, 7.18, 8.0, 0.3,
        size=9, color=MUTED,
    )
    _add_text(
        slide, "Plan. Fly. Land.", 12.0, 7.18, 1.0, 0.3,
        size=9, color=MUTED, align=PP_ALIGN.RIGHT,
    )


def build_slide_one(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_rect(slide, 0, 0, prs.slide_width / 914400, prs.slide_height / 914400, PAPER)
    _slide_chrome(slide, slide_num=1, total=2, eyebrow="THE PIVOT")

    # Title
    _add_text(
        slide, "ClickShop  →  Meridian", 0.5, 0.85, 12.5, 0.7,
        size=36, bold=True, color=INK,
    )
    _add_text(
        slide,
        "Why we pivoted the DAT309 demo",
        0.5, 1.55, 12.5, 0.4,
        size=16, color=INK_SOFT,
    )

    # Lede
    _add_rule(slide, 0.5, 2.1, 1.5, ACCENT)
    _add_text(
        slide,
        "The abstract promises 5 things.  ClickShop showed 2.  Meridian shows 5.",
        0.5, 2.25, 12.5, 0.5,
        size=18, bold=True, color=INK,
    )

    # Coverage table
    table_top = 3.05
    row_h = 0.46
    col_x = [0.5, 8.7, 10.85]   # claim | clickshop | meridian
    col_w = [8.0, 2.0, 2.0]

    # Header
    header_bg = _add_rect(slide, 0.5, table_top, 12.5, row_h, ACCENT_SOFT)
    _add_text(slide, "Abstract claim", col_x[0] + 0.18, table_top + 0.10, col_w[0] - 0.2, row_h,
              size=11, bold=True, color=ACCENT)
    _add_text(slide, "ClickShop", col_x[1], table_top + 0.10, col_w[1], row_h,
              size=11, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    _add_text(slide, "Meridian", col_x[2], table_top + 0.10, col_w[2], row_h,
              size=11, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    rows = [
        ("Aurora + MCP",                                   True,  True),
        ("Bedrock agents on live relational data",         True,  True),
        ("MCP servers for contextual memory",              False, True),
        ("Securely connect LLM agents to Aurora",          False, True),
        ("AgentCore + LangGraph + Strands workflows",      False, True),
    ]

    for i, (claim, click, mer) in enumerate(rows):
        y = table_top + row_h * (i + 1)
        if i % 2 == 0:
            _add_rect(slide, 0.5, y, 12.5, row_h, RGBColor(0xF1, 0xEC, 0xE2))
        _add_text(slide, claim, col_x[0] + 0.18, y + 0.10, col_w[0] - 0.2, row_h,
                  size=12, color=INK)
        _add_text(slide, "✓" if click else "—", col_x[1], y + 0.07, col_w[1], row_h,
                  size=18, bold=True,
                  color=GREEN if click else MUTED, align=PP_ALIGN.CENTER)
        _add_text(slide, "✓" if mer else "—", col_x[2], y + 0.07, col_w[2], row_h,
                  size=18, bold=True,
                  color=GREEN if mer else MUTED, align=PP_ALIGN.CENTER)

    # Closing kicker
    kicker_top = table_top + row_h * (len(rows) + 1) + 0.25
    _add_rule(slide, 0.5, kicker_top, 1.5, ACCENT)
    _add_text(
        slide,
        "Travel is also a better foil for the agentic story —",
        0.5, kicker_top + 0.15, 12.5, 0.4,
        size=14, color=INK_SOFT,
    )
    _add_text(
        slide,
        "a returning traveler with preferences, allergies, and prior trips  >  a shopping cart.",
        0.5, kicker_top + 0.55, 12.5, 0.4,
        size=14, bold=True, color=INK,
    )


def build_slide_two(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, prs.slide_width / 914400, prs.slide_height / 914400, PAPER)
    _slide_chrome(slide, slide_num=2, total=2, eyebrow="THE LADDER")

    _add_text(slide, "Meridian — the five-phase ladder", 0.5, 0.85, 12.5, 0.7,
              size=32, bold=True, color=INK)
    _add_text(
        slide,
        "Each phase adds exactly one capability on the same Aurora catalog.",
        0.5, 1.55, 12.5, 0.4,
        size=14, color=INK_SOFT,
    )

    phases = [
        ("1", "SQL",
         "RDS Data API → trip_packages",
         "“City trips under $2,000” works  ·  “Romantic wine-country week” returns 0"),
        ("2", "MCP",
         "postgres-mcp-server (public)",
         "Same intent gap — interface change, not intelligence"),
        ("3", "Retrieval",
         "Cohere Embed v4 (1024d) + pgvector + tsvector  ·  Strands supervisor",
         "“Romantic week in Europe” now ranks correctly"),
        ("4", "Memory",
         "ConciergeOrchestrator  ·  Aurora memory + RLS + audit  ·  AgentCore Memory & Identity",
         "Knows Alex Morgan, Tokyo Oct 12–19, shellfish allergy"),
        ("5", "Orchestration",
         "LangGraph StateGraph  ·  PostgresSaver checkpoints in Aurora",
         "classify → search / availability / recall → synthesize"),
    ]

    top = 2.15
    row_h = 0.78
    for i, (num, title, capability, hook) in enumerate(phases):
        y = top + i * row_h
        # Number chip
        _add_rect(slide, 0.5, y + 0.06, 0.55, 0.55, PHASE_COLORS[i])
        _add_text(slide, num, 0.5, y + 0.13, 0.55, 0.55,
                  size=22, bold=True, color=PAPER, align=PP_ALIGN.CENTER)
        # Phase title
        _add_text(slide, title, 1.25, y + 0.04, 3.0, 0.4,
                  size=18, bold=True, color=INK)
        # Capability
        _add_text(slide, capability, 1.25, y + 0.36, 11.6, 0.32,
                  size=11, color=INK_SOFT)
        # Demo hook
        _add_text(slide, hook, 4.5, y + 0.04, 8.4, 0.32,
                  size=11, color=PHASE_COLORS[i], bold=True, align=PP_ALIGN.RIGHT)
        if i < len(phases) - 1:
            _add_rule(slide, 1.25, y + row_h - 0.04, 11.5)

    # Footer band: framing + status
    band_top = top + len(phases) * row_h + 0.2
    _add_rect(slide, 0.5, band_top, 12.5, 0.85, ACCENT_SOFT)
    _add_text(
        slide,
        "Strands routes tools in Phases 3–4.  LangGraph owns control flow in Phase 5.",
        0.7, band_top + 0.10, 12.0, 0.32,
        size=12, bold=True, color=ACCENT,
    )
    _add_text(
        slide,
        "Custom meridian-memory MCP server lives alongside the public postgres-mcp-server.",
        0.7, band_top + 0.40, 12.0, 0.32,
        size=11, color=INK_SOFT,
    )


def main() -> None:
    repo = Path(__file__).resolve().parents[1]
    out_dir = repo / "docs"
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / "meridian-pivot.pptx"

    prs = Presentation()
    # 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_slide_one(prs)
    build_slide_two(prs)

    prs.save(out_path)
    print(f"wrote {out_path.relative_to(repo)}")


if __name__ == "__main__":
    main()
